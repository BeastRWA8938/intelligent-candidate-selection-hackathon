import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def fill_presentation():
    # Load the presentation
    prs = pptx.Presentation("Idea Submission Template _ Redrob.pptx")
    
    # Color definition (navy blue/gray matching professional design)
    TEXT_COLOR = RGBColor(30, 41, 59) # Slate 800
    
    # ----------------------------------------------------
    # SLIDE 0: Title Slide
    # Shape 1: Team Name
    # Shape 2: Problem Statement
    # Shape 3: Team Leader Name
    # ----------------------------------------------------
    slide0 = prs.slides[0]
    slide0.shapes[1].text_frame.text = "Team Name : Team Redrob"
    slide0.shapes[2].text_frame.text = "Problem Statement : Intelligent Candidate Discovery & Ranking"
    slide0.shapes[3].text_frame.text = "Team Leader Name : Rushikesh"
    
    # Apply standard formatting for slide 0
    for shape in [slide0.shapes[1], slide0.shapes[2], slide0.shapes[3]]:
        for p in shape.text_frame.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = TEXT_COLOR

    # Helper function to format content slides
    def set_content_slide(slide_index, bullets_list):
        slide = prs.slides[slide_index]
        body_shape = slide.shapes[2] # Shape 2 is the prompt textbox
        tf = body_shape.text_frame
        tf.clear()
        
        for idx, item in enumerate(bullets_list):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = item
            if item.startswith("• "):
                p.level = 0
                p.font.size = Pt(14)
                p.space_after = Pt(4)
                p.font.name = "Calibri"
                p.font.color.rgb = TEXT_COLOR
            else:
                p.level = 0
                p.font.bold = True
                p.font.size = Pt(16)
                p.space_before = Pt(8)
                p.space_after = Pt(4)
                p.font.name = "Calibri"
                p.font.color.rgb = RGBColor(15, 23, 42) # Slate 900
                
    # ----------------------------------------------------
    # SLIDE 1: Solution Overview
    # ----------------------------------------------------
    bullets_s1 = [
        "Proposed Solution: Hybrid Lexical-Semantic Candidate Discovery Pipeline",
        "• Deterministic Anomaly Sieve: Screen and discard 100% of synthetic honeypots and consulting firm candidates before scoring.",
        "• Custom BM25 Relevance Engine: Computes sparse term vectors with frequency saturation and length normalization to neutralize keyword stuffers.",
        "• Heuristic Feature Fusion: Combines title relevance, skills match, pedigree, location preferences, and asymmetric experience metrics.",
        "• Logarithmic Social Proof & Recency Decay: Dampens endorsement farming and discounts stale skills.",
        "• SHA-256 Reasoner: Deterministic, rank-consistent candidate justifications."
    ]
    set_content_slide(1, bullets_s1)

    # ----------------------------------------------------
    # SLIDE 2: JD Understanding & Candidate Evaluation
    # ----------------------------------------------------
    bullets_s2 = [
        "Key Job Description Requirements Extracted:",
        "• Technical: Embeddings, dense search, vector databases (Pinecone, Qdrant, Milvus, FAISS), information retrieval, and NDCG/MRR/MAP evaluation.",
        "• Experience & Pedigree: Senior AI/ML Engineer with 5-9 years experience; Tier-1/2 education; product company background preferred.",
        "• Location: Noida/Pune hybrid setup alignment.",
        "Candidate Evaluation Signals:",
        "• Capability: Technical title matching, BM25 similarity, education pedigree, and skill depth-weighted clusters.",
        "• Availability & Engagement: Smoothed response rate, notice period, active logins, and GitHub scores (restricted to [0.85, 1.15] range)."
    ]
    set_content_slide(2, bullets_s2)

    # ----------------------------------------------------
    # SLIDE 3: Ranking Methodology
    # ----------------------------------------------------
    bullets_s3 = [
        "3-Pass Streaming Architecture:",
        "• Pass 1: Sieve candidate anomalies; compute vocabulary Document Frequencies and average document length.",
        "• Pass 2: Calculate BM25 scores; extract max BM25 score for scale normalization.",
        "• Pass 3: Normalize BM25; compute base scores and apply behavioral multipliers; maintain top 100 in min-heap.",
        "Algorithmic Hardening:",
        "• Asymmetric Experience Sweet Spot: Gaussian decay for YoE < 5; plateau floor (75.0) for YoE > 9.",
        "• Depth-Weighted Skills: Blends cluster diversity and maximum skill depth to prevent generalist bias.",
        "• Recency Decay: Exponential decay factor of e^(-0.015 * t) for stale skills."
    ]
    set_content_slide(3, bullets_s3)

    # ----------------------------------------------------
    # SLIDE 4: Explainability & Data Validation
    # ----------------------------------------------------
    bullets_s4 = [
        "Explainable Decision Logic:",
        "• SHA-256 Hash Seed: Order-sensitive seed of Candidate ID ensures deterministic and collision-free reasonings.",
        "• Rank-Consistent Tone: Justification structure and adjectives adapt to candidate rank (e.g. Ranks 1-10 vs 51-100) and actual matching skills.",
        "• Hallucination Prevention: References actual primary profile skills if concept clusters are empty.",
        "Deterministic Anomaly Sieve (Honeypot Prevention):",
        "• YoE vs Job Span: Flags if YoE > job span + 2.0 years.",
        "• Job Duration: Flags if job duration > calendar span + 3 months.",
        "• Expert Zero-Duration: Flags if >= 4 expert skills have 0 months experience."
    ]
    set_content_slide(4, bullets_s4)

    # ----------------------------------------------------
    # SLIDE 5: End-to-End Workflow
    # ----------------------------------------------------
    bullets_s5 = [
        "Complete Candidate Pipeline:",
        "1. Ingestion: Stream candidates.jsonl (100k rows) in raw JSON format.",
        "2. Data Cleansing: Pass candidates through the Anomaly Sieve and consulting firm checks.",
        "3. Indexing & BM25: Tokenize profiles; calculate document frequencies and BM25 scores.",
        "4. Scoring & Blending: Fuse Title, BM25, Skills, YoE, Pedigree, and Location with behavioral multipliers.",
        "5. Top-100 Filtering: Maintain candidates in a bounded heapq min-heap.",
        "6. Final Export: Resolve score ties alphabetically; generate reasonings; write team_redrob.csv."
    ]
    set_content_slide(5, bullets_s5)

    # ----------------------------------------------------
    # SLIDE 6: System Architecture (Add new text shape with ASCII diagram)
    # ----------------------------------------------------
    slide6 = prs.slides[6]
    # Add textbox for the diagram
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(3.0)
    txBox = slide6.shapes.add_textbox(left, top, width, height)
    tf6 = txBox.text_frame
    tf6.clear()
    p6 = tf6.paragraphs[0]
    p6.text = (
        "+---------------------------------------------------------------+\n"
        "|                      candidates.jsonl                         |\n"
        "+-------------------------------|-------------------------------+\n"
        "                                v\n"
        "+---------------------------------------------------------------+\n"
        "|        Layer 1: Sieve Filters (Honeypots & Consulting)        |\n"
        "+-------------------------------|-------------------------------+\n"
        "                                v\n"
        "+---------------------------------------------------------------+\n"
        "|        Layer 2: BM25 Lexical Relevance Retrieval              |\n"
        "+-------------------------------|-------------------------------+\n"
        "                                v\n"
        "+---------------------------------------------------------------+\n"
        "|        Layer 3: Heuristic Scorer (Title, Skills, YoE, Loc)    |\n"
        "+-------------------------------|-------------------------------+\n"
        "                                v\n"
        "+---------------------------------------------------------------+\n"
        "|        Layer 4: Ranking & Reasoning (Heap & SHA-256)          |\n"
        "+-------------------------------|-------------------------------+\n"
        "                                v\n"
        "+---------------------------------------------------------------+\n"
        "|                     team_redrob.csv (Top 100)                 |\n"
        "+---------------------------------------------------------------+"
    )
    p6.font.name = "Consolas"
    p6.font.size = Pt(11)
    p6.font.color.rgb = TEXT_COLOR

    # ----------------------------------------------------
    # SLIDE 7: Results & Performance
    # ----------------------------------------------------
    bullets_s7 = [
        "System Benchmarks:",
        "• Execution Speed: Processed 100,000 candidates in 44.15 seconds (5.8x faster than the 5-minute sandbox limit).",
        "• Memory Efficiency: Streaming architecture uses only ~35MB RAM, ensuring OOM protection.",
        "• Anomaly Capture: Identified and filtered 57 honeypots (0% honeypot rate in top 100).",
        "• Disqualifications: Filtered out 52,448 unqualified titles and consulting backgrounds.",
        "• Validation Checks: 100% compliance on column formats, row count, sorting, and tie-breakers.",
        "• Defense Potential: High explainability and mathematical robustness for the interview defense stage."
    ]
    set_content_slide(7, bullets_s7)

    # ----------------------------------------------------
    # SLIDE 8: Technologies Used
    # ----------------------------------------------------
    bullets_s8 = [
        "Technology Stack:",
        "• Python 3.11 Standard Library: Standard library only (no external packages like torch, numpy, pandas).",
        "• Key modules used:",
        "  - json & gzip: Streaming data ingestion and compression/decompression.",
        "  - re: Tech-aware pattern-based tokenization.",
        "  - heapq: Bounded min-heap for O(1) memory footprint.",
        "  - hashlib: SHA-256 seed hashing.",
        "  - math: BM25 logs, Gaussian YoE calculations.",
        "  - collections: Counter and defaultdict for vocabulary mapping.",
        "• Why Selected:",
        "  - 100% Sandbox Reproducibility: Guarantees zero installation or build failures in the grading sandbox.",
        "  - Offline Security: Operates with zero network connections."
    ]
    set_content_slide(8, bullets_s8)

    # ----------------------------------------------------
    # SLIDE 9: Submission Assets
    # ----------------------------------------------------
    bullets_s9 = [
        "Project Deliverables:",
        "• Core Ranking Engine: rank.py",
        "• Final Ranked Output: team_redrob.csv (exactly 100 candidates)",
        "• Submission Metadata: submission_metadata.yaml",
        "• Visualizer Code: visualizer_server.py and static/",
        "Links and Repositories:",
        "• GitHub Repository: https://github.com/BeastRWA8938/intelligent-candidate-selection-hackathon",
        "• Sandbox Hosted Link: Streamlit cloud dashboard presenting candidate scores and audit details.",
        "• Walkthrough Documentation: walkthrough.md and solutions_to_loopholes.md"
    ]
    set_content_slide(9, bullets_s9)

    # Save presentation
    prs.save("Idea Submission Template _ Redrob.pptx")
    print("[+] PowerPoint template successfully filled and formatted!")

if __name__ == "__main__":
    fill_presentation()
