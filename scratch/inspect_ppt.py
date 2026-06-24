import pptx

prs = pptx.Presentation("Idea Submission Template _ Redrob.pptx")
for idx, slide in enumerate(prs.slides):
    print(f"=== SLIDE {idx} ===")
    for s_idx, shape in enumerate(slide.shapes):
        has_tf = hasattr(shape, "text_frame") and shape.text_frame is not None
        text = shape.text_frame.text.strip() if has_tf else ""
        if text:
            # Replace newlines for clean printing
            clean_text = text.replace("\n", " | ")
            print(f"  Shape {s_idx} ({shape.name}, Type={shape.shape_type}): \"{clean_text}\"")
        else:
            print(f"  Shape {s_idx} ({shape.name}, Type={shape.shape_type}): [No Text]")
